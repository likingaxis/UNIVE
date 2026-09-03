#!/usr/bin/env python3
"""
========================================================================================
BUILD EDITABLE POWERPOINT DECK: DONNE, GAMING E INCLUSIONE
========================================================================================
Generates a 100% editable, native vector PowerPoint (.pptx) presentation with:
- 16:9 widescreen layout (13.333" x 7.5")
- Native text boxes, fonts (Segoe UI / Consolas)
- Native rectangular cards with subtle borders & accent lines
- Editable huge statistical callouts (76%, 48%, 23-25%, -24%, 68% vs 38%, 33%)
- Native overlapping progress bars (track + fill)
- Faithful presentation of all 9 approved slides
========================================================================================
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --------------------------------------------------------------------------------------
# COLOR PALETTE (RGB)
# --------------------------------------------------------------------------------------
COLOR_BG            = RGBColor(7, 17, 31)       # #07111F - Main background
COLOR_SURFACE       = RGBColor(15, 33, 56)      # #0F2138 - Card background
COLOR_SURFACE_ALT   = RGBColor(18, 42, 70)      # #122A46 - Highlighted surface
COLOR_SURFACE_SUB   = RGBColor(11, 24, 43)      # #0B182B - Deeper nested card
COLOR_BORDER_SUBTLE = RGBColor(35, 58, 88)      # #233A58 - Subtle card border
COLOR_BORDER_MUTED  = RGBColor(48, 75, 110)     # #304B6E - Medium card border

COLOR_TEXT_MAIN     = RGBColor(244, 247, 251)   # #F4F7FB - Primary text
COLOR_TEXT_MUTED    = RGBColor(169, 184, 204)   # #A9B8CC - Muted text / descriptions
COLOR_WHITE         = RGBColor(255, 255, 255)   # #FFFFFF

COLOR_CYAN          = RGBColor(55, 214, 232)    # #37D6E8 - Tech Cyan (Speaker 1)
COLOR_VIOLET        = RGBColor(167, 139, 250)   # #A78BFA - Inclusion Violet (Speaker 2)
COLOR_CORAL         = RGBColor(255, 107, 107)   # #FF6B6B - Alert Coral
COLOR_MINT          = RGBColor(94, 230, 168)    # #5EE6A8 - Positive Mint (Speaker 3)

COLOR_TRACK         = RGBColor(24, 45, 72)      # Progress bar background track

# --------------------------------------------------------------------------------------
# TYPOGRAPHY CONSTANTS
# --------------------------------------------------------------------------------------
FONT_TITLE = "Segoe UI"
FONT_BODY  = "Segoe UI"
FONT_MONO  = "Consolas"

# --------------------------------------------------------------------------------------
# DECK BUILDER CLASS
# --------------------------------------------------------------------------------------
class EditableDeckBuilder:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.blank_layout = self.prs.slide_layouts[6]  # blank layout

    def create_slide(self):
        """Creates a slide with dark background #07111F."""
        slide = self.prs.slides.add_slide(self.blank_layout)
        # Native slide background fill
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

        # Underlying full rectangle guarantee for cross-platform compatibility
        bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = COLOR_BG
        bg_shape.line.fill.background()
        return slide

    def add_speaker_badge(self, slide, speaker_text, accent_color):
        """Creates top-right speaker badge with colored outline."""
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(9.8), Inches(0.42), Inches(2.733), Inches(0.38)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = COLOR_SURFACE
        badge.line.color.rgb = accent_color
        badge.line.width = Pt(1.2)

        tf = badge.text_frame
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = speaker_text
        p.font.name = FONT_MONO
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = accent_color

    def add_header(self, slide, kicker_text, kicker_color, title_text, subtitle_text, speaker_text=None, speaker_color=None):
        """Adds consistent kicker, title H2, subtitle and speaker badge."""
        # Speaker Badge
        if speaker_text:
            self.add_speaker_badge(slide, speaker_text, speaker_color or kicker_color)

        # Kicker (Mono, All-caps)
        kicker_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.44), Inches(8.8), Inches(0.32))
        tf_k = kicker_box.text_frame
        tf_k.word_wrap = False
        tf_k.margin_left = tf_k.margin_right = tf_k.margin_top = tf_k.margin_bottom = 0
        pk = tf_k.paragraphs[0]
        pk.text = kicker_text
        pk.font.name = FONT_MONO
        pk.font.size = Pt(11.5)
        pk.font.bold = True
        pk.font.color.rgb = kicker_color

        # Title H2
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.82), Inches(11.7), Inches(0.75))
        tf_t = title_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_right = tf_t.margin_top = tf_t.margin_bottom = 0
        pt = tf_t.paragraphs[0]
        pt.text = title_text
        pt.font.name = FONT_TITLE
        pt.font.size = Pt(29)
        pt.font.bold = True
        pt.font.color.rgb = COLOR_TEXT_MAIN

        # Subtitle
        if subtitle_text:
            sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.60), Inches(11.7), Inches(0.45))
            tf_s = sub_box.text_frame
            tf_s.word_wrap = True
            tf_s.margin_left = tf_s.margin_right = tf_s.margin_top = tf_s.margin_bottom = 0
            ps = tf_s.paragraphs[0]
            ps.text = subtitle_text
            ps.font.name = FONT_BODY
            ps.font.size = Pt(13.5)
            ps.font.color.rgb = COLOR_TEXT_MUTED

    def add_footer(self, slide, slide_num, total_slides=9):
        """Adds academic footer and slide pagination."""
        # Footer Left
        foot_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.92), Inches(8.0), Inches(0.3))
        tf_f = foot_box.text_frame
        tf_f.margin_left = tf_f.margin_right = tf_f.margin_top = tf_f.margin_bottom = 0
        pf = tf_f.paragraphs[0]
        pf.text = "Gender & Inclusion · Università di Roma Tor Vergata"
        pf.font.name = FONT_MONO
        pf.font.size = Pt(10)
        pf.font.color.rgb = RGBColor(120, 140, 168)

        # Slide Number Right
        num_box = slide.shapes.add_textbox(Inches(10.5), Inches(6.92), Inches(2.033), Inches(0.3))
        tf_n = num_box.text_frame
        tf_n.margin_left = tf_n.margin_right = tf_n.margin_top = tf_n.margin_bottom = 0
        pn = tf_n.paragraphs[0]
        pn.alignment = PP_ALIGN.RIGHT
        pn.text = f"{slide_num:02d} / {total_slides:02d}"
        pn.font.name = FONT_MONO
        pn.font.size = Pt(10)
        pn.font.color.rgb = RGBColor(120, 140, 168)

    def add_card(self, slide, left, top, width, height, bg_color=COLOR_SURFACE, border_color=COLOR_BORDER_SUBTLE, left_accent=None, left_accent_width=Inches(0.065)):
        """Creates an editable card with optional left accent stripe."""
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        if border_color:
            card.line.color.rgb = border_color
            card.line.width = Pt(1)
        else:
            card.line.fill.background()

        # Left colored accent line
        if left_accent:
            accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + Inches(0.05), left_accent_width, height - Inches(0.1))
            accent_bar.fill.solid()
            accent_bar.fill.fore_color.rgb = left_accent
            accent_bar.line.fill.background()

        return card

    def add_quote_bar(self, slide, left, top, width, height, regular_text, bold_highlight, accent_color=COLOR_CYAN):
        """Creates high-impact quote bar with vertical accent."""
        # Vertical accent line
        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + Inches(0.04), Inches(0.055), height - Inches(0.08))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent_color
        accent_bar.line.fill.background()

        # Text Frame
        tb = slide.shapes.add_textbox(left + Inches(0.18), top, width - Inches(0.18), height)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]

        r1 = p.add_run()
        r1.text = regular_text
        r1.font.name = FONT_BODY
        r1.font.size = Pt(17)
        r1.font.color.rgb = COLOR_TEXT_MAIN

        r2 = p.add_run()
        r2.text = bold_highlight
        r2.font.name = FONT_BODY
        r2.font.size = Pt(17)
        r2.font.bold = True
        r2.font.color.rgb = accent_color

    def add_progress_bar(self, slide, left, top, width, height, percent, fill_color, track_color=COLOR_TRACK):
        """Adds a native progress bar with track and fill rectangles."""
        # Track
        track = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        track.fill.solid()
        track.fill.fore_color.rgb = track_color
        track.line.fill.background()

        # Fill
        fill_w = max(width * percent, Inches(0.1))
        fill_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, fill_w, height)
        fill_shape.fill.solid()
        fill_shape.fill.fore_color.rgb = fill_color
        fill_shape.line.fill.background()

    # ----------------------------------------------------------------------------------
    # SLIDE 1: TITLE & COVER SLIDE
    # ----------------------------------------------------------------------------------
    def build_slide_1(self):
        slide = self.create_slide()

        # Academic Kicker Tag (Rounded Pill)
        kicker_tag = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(0.7), Inches(5.8), Inches(0.38)
        )
        kicker_tag.fill.solid()
        kicker_tag.fill.fore_color.rgb = COLOR_SURFACE
        kicker_tag.line.color.rgb = COLOR_CYAN
        kicker_tag.line.width = Pt(1.2)
        tf_kt = kicker_tag.text_frame
        tf_kt.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf_kt.margin_left = Inches(0.15)
        p_kt = tf_kt.paragraphs[0]
        p_kt.text = "● UNIVERSITÀ DI ROMA TOR VERGATA · INFORMATICA"
        p_kt.font.name = FONT_MONO
        p_kt.font.size = Pt(11)
        p_kt.font.bold = True
        p_kt.font.color.rgb = COLOR_CYAN

        # Big Main Title H1
        h1_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.28), Inches(7.5), Inches(1.8))
        tf_h1 = h1_box.text_frame
        tf_h1.word_wrap = True
        tf_h1.margin_left = tf_h1.margin_right = tf_h1.margin_top = tf_h1.margin_bottom = 0
        p_h1_1 = tf_h1.paragraphs[0]
        p_h1_1.text = "Donne, gaming"
        p_h1_1.font.name = FONT_TITLE
        p_h1_1.font.size = Pt(48)
        p_h1_1.font.bold = True
        p_h1_1.font.color.rgb = COLOR_TEXT_MAIN

        p_h1_2 = tf_h1.add_paragraph()
        r_inc = p_h1_2.add_run()
        r_inc.text = "e inclusione"
        r_inc.font.name = FONT_TITLE
        r_inc.font.size = Pt(48)
        r_inc.font.bold = True
        r_inc.font.color.rgb = COLOR_VIOLET

        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.18), Inches(7.2), Inches(0.9))
        tf_sub = sub_box.text_frame
        tf_sub.word_wrap = True
        tf_sub.margin_left = tf_sub.margin_right = tf_sub.margin_top = tf_sub.margin_bottom = 0
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = "Donne e cultura videoludica: radici storiche, dati di settore e considerazioni sui futuri sviluppi tra codice e società"
        p_sub.font.name = FONT_BODY
        p_sub.font.size = Pt(15.5)
        p_sub.font.color.rgb = COLOR_TEXT_MAIN

        # Team Pills (Speakers 01, 02, 03)
        team_y = Inches(4.32)
        members = [
            ("01", "Luca Gugliotta", "0342634", COLOR_CYAN),
            ("02", "Valerio Bernardi", "0349538", COLOR_VIOLET),
            ("03", "Samuele De Santis", "0348324", COLOR_MINT),
        ]
        for badge_no, name, matricola, color in members:
            pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), team_y, Inches(5.8), Inches(0.44))
            pill.fill.solid()
            pill.fill.fore_color.rgb = COLOR_SURFACE
            pill.line.color.rgb = COLOR_BORDER_SUBTLE
            pill.line.width = Pt(1)

            # Left accent
            accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), team_y + Inches(0.04), Inches(0.065), Inches(0.36))
            accent.fill.solid()
            accent.fill.fore_color.rgb = color
            accent.line.fill.background()

            tf_p = pill.text_frame
            tf_p.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf_p.margin_left = Inches(0.2)
            p_p = tf_p.paragraphs[0]

            r_num = p_p.add_run()
            r_num.text = f"{badge_no}   "
            r_num.font.name = FONT_MONO
            r_num.font.bold = True
            r_num.font.size = Pt(12)
            r_num.font.color.rgb = color

            r_name = p_p.add_run()
            r_name.text = f"{name}   "
            r_name.font.name = FONT_BODY
            r_name.font.bold = True
            r_name.font.size = Pt(13)
            r_name.font.color.rgb = COLOR_TEXT_MAIN

            r_matr = p_p.add_run()
            r_matr.text = f"Matr. {matricola}"
            r_matr.font.name = FONT_MONO
            r_matr.font.size = Pt(11.5)
            r_matr.font.color.rgb = COLOR_TEXT_MUTED

            team_y += Inches(0.54)

        # Session Meta
        meta_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.15), Inches(7.0), Inches(0.35))
        tf_m = meta_box.text_frame
        p_m = tf_m.paragraphs[0]
        p_m.text = "GENDER & INCLUSION · A.A. 2025–2026 · 4/09/2026"
        p_m.font.name = FONT_MONO
        p_m.font.size = Pt(11.5)
        p_m.font.color.rgb = COLOR_TEXT_MUTED

        # Right Side Visual Card (Academic & Gaming HUD Card)
        right_card = self.add_card(
            slide, Inches(8.4), Inches(0.95), Inches(4.133), Inches(5.5),
            bg_color=COLOR_SURFACE, border_color=COLOR_CYAN, left_accent=COLOR_CYAN
        )
        tb_rc = slide.shapes.add_textbox(Inches(8.65), Inches(1.15), Inches(3.65), Inches(5.1))
        tf_rc = tb_rc.text_frame
        tf_rc.word_wrap = True

        p1 = tf_rc.paragraphs[0]
        p1.text = "RESEARCH FRAMEWORK"
        p1.font.name = FONT_MONO
        p1.font.bold = True
        p1.font.size = Pt(12)
        p1.font.color.rgb = COLOR_CYAN

        items = [
            ("OBIETTIVO PROGETTO", "Analisi interdisciplinare del medium videoludico: demografia, tossicità online, evoluzione narrativa e disparità occupazionale nelle carriere STEM."),
            ("DATASET & FONTI", "ESA 2025 · ADL Hate is No Game 2024 · IIDEA 2025 · GDC State of Industry 2025 · VGE 2024."),
            ("RELAZIONE COLLEGIALE", "Slide 01-03: Luca Gugliotta\nSlide 04-05: Valerio Bernardi\nSlide 06-09: Samuele De Santis"),
            ("STATUS METRICHE", "PARITÀ CONSUMO: 48% F / 52% M\nSVILUPPO TECH: 20-23% DONNE\nONLINE HARASSMENT: 76% COMP.")
        ]
        for title, desc in items:
            p_t = tf_rc.add_paragraph()
            p_t.space_before = Pt(10)
            p_t.text = f"■ {title}"
            p_t.font.name = FONT_MONO
            p_t.font.bold = True
            p_t.font.size = Pt(10.5)
            p_t.font.color.rgb = COLOR_VIOLET

            p_d = tf_rc.add_paragraph()
            p_d.text = desc
            p_d.font.name = FONT_BODY
            p_d.font.size = Pt(10.5)
            p_d.font.color.rgb = COLOR_TEXT_MUTED

    # ----------------------------------------------------------------------------------
    # SLIDE 2: PERCHÉ PARLARNE DA INFORMATICI? (Speaker 1)
    # ----------------------------------------------------------------------------------
    def build_slide_2(self):
        slide = self.create_slide()
        self.add_header(
            slide,
            kicker_text="01 / LA PROSPETTIVA INFORMATICA",
            kicker_color=COLOR_CYAN,
            title_text="Perché parlarne da informatici?",
            subtitle_text="I videogiochi sono spazi sociali: le scelte di codice e design determinano chi può partecipare in sicurezza.",
            speaker_text="01 · LUCA GUGLIOTTA",
            speaker_color=COLOR_CYAN
        )
        self.add_footer(slide, 2)

        # 3 Pillar Cards
        cols = [
            (
                "01", "CHI GIOCA", COLOR_CYAN,
                "Voice chat & matchmaking",
                "Tossicità, sicurezza e moderazione",
                "Gestione dell'anonimato e sistemi anti-molestia nell'esperienza multiplayer."
            ),
            (
                "02", "COSA VEDIAMO", COLOR_VIOLET,
                "Narrazione e stereotipi",
                "Evoluzione del character design",
                "Dalla figura passiva del passato all'agentività e all'identità modulare attuale."
            ),
            (
                "03", "CHI SVILUPPA", COLOR_MINT,
                "Donne nella carriera gaming",
                "Ruoli tecnici, leadership e divario",
                "Composizione dei team di produzione, retribuzioni e cultura dell'ingegneria software."
            )
        ]

        left_start = Inches(0.8)
        card_w = Inches(3.644)
        card_h = Inches(3.72)
        gap = Inches(0.4)

        for i, (num, title, color, pt1, pt2, detail) in enumerate(cols):
            x = left_start + i * (card_w + gap)
            y = Inches(2.28)

            card = self.add_card(slide, x, y, card_w, card_h, bg_color=COLOR_SURFACE, border_color=color, left_accent=color)

            tb = slide.shapes.add_textbox(x + Inches(0.22), y + Inches(0.18), card_w - Inches(0.4), card_h - Inches(0.3))
            tf = tb.text_frame
            tf.word_wrap = True

            # Number badge
            p_num = tf.paragraphs[0]
            p_num.text = num
            p_num.font.name = FONT_MONO
            p_num.font.size = Pt(46)
            p_num.font.bold = True
            p_num.font.color.rgb = color

            # Pillar Title
            p_title = tf.add_paragraph()
            p_title.text = title
            p_title.font.name = FONT_TITLE
            p_title.font.size = Pt(20)
            p_title.font.bold = True
            p_title.font.color.rgb = COLOR_TEXT_MAIN
            p_title.space_after = Pt(12)

            # Bullet 1
            p_b1 = tf.add_paragraph()
            p_b1.text = f"• {pt1}"
            p_b1.font.name = FONT_BODY
            p_b1.font.bold = True
            p_b1.font.size = Pt(13)
            p_b1.font.color.rgb = COLOR_TEXT_MAIN

            # Bullet 2
            p_b2 = tf.add_paragraph()
            p_b2.space_before = Pt(4)
            p_b2.text = f"• {pt2}"
            p_b2.font.name = FONT_BODY
            p_b2.font.bold = True
            p_b2.font.size = Pt(13)
            p_b2.font.color.rgb = COLOR_TEXT_MAIN

            # Detail paragraph
            p_det = tf.add_paragraph()
            p_det.space_before = Pt(12)
            p_det.text = detail
            p_det.font.name = FONT_BODY
            p_det.font.size = Pt(11.5)
            p_det.font.color.rgb = COLOR_TEXT_MUTED

        # Quote Bar at Bottom
        self.add_quote_bar(
            slide,
            left=Inches(0.8), top=Inches(6.16), width=Inches(11.733), height=Inches(0.6),
            regular_text="Il modo in cui programmiamo uno spazio digitale determina ",
            bold_highlight="chi potrà viverlo.",
            accent_color=COLOR_CYAN
        )

    # ----------------------------------------------------------------------------------
    # SLIDE 3: DEMOGRAFIA E MERCATO (Speaker 1)
    # ----------------------------------------------------------------------------------
    def build_slide_3(self):
        slide = self.create_slide()
        self.add_header(
            slide,
            kicker_text="02 / DEMOGRAFIA E MERCATO",
            kicker_color=COLOR_CYAN,
            title_text="La realtà è quasi paritaria. La percezione no.",
            subtitle_text="Il paradosso percettivo: la partecipazione reale sfiora la parità, ma lo stereotipo del gamer resta maschile.",
            speaker_text="01 · LUCA GUGLIOTTA",
            speaker_color=COLOR_CYAN
        )
        self.add_footer(slide, 3)

        # 2 Main Columns
        col_w = Inches(5.666)
        col_h = Inches(3.72)
        y = Inches(2.28)

        # Left Column: Global Audience (ESA 2025)
        left_x = Inches(0.8)
        self.add_card(slide, left_x, y, col_w, col_h, bg_color=COLOR_SURFACE, border_color=COLOR_BORDER_SUBTLE, left_accent=COLOR_CYAN)

        tb_l = slide.shapes.add_textbox(left_x + Inches(0.25), y + Inches(0.18), col_w - Inches(0.5), Inches(1.1))
        tf_l = tb_l.text_frame
        tf_l.word_wrap = True

        p_lk = tf_l.paragraphs[0]
        p_lk.text = "PUBBLICO GLOBALE · ESA 2025"
        p_lk.font.name = FONT_MONO
        p_lk.font.size = Pt(11.5)
        p_lk.font.bold = True
        p_lk.font.color.rgb = COLOR_CYAN

        # Parity Percentage Labels
        p_lbl = tf_l.add_paragraph()
        p_lbl.space_before = Pt(8)
        r_w = p_lbl.add_run()
        r_w.text = "48% DONNE        "
        r_w.font.name = FONT_MONO
        r_w.font.size = Pt(22)
        r_w.font.bold = True
        r_w.font.color.rgb = COLOR_CYAN

        r_m = p_lbl.add_run()
        r_m.text = "52% UOMINI"
        r_m.font.name = FONT_MONO
        r_m.font.size = Pt(22)
        r_m.font.bold = True
        r_m.font.color.rgb = COLOR_VIOLET

        # Native Split Bar (48% Cyan + 52% Violet + Center Marker)
        bar_x = left_x + Inches(0.25)
        bar_y = y + Inches(1.42)
        bar_w = col_w - Inches(0.5)
        bar_h = Inches(0.32)

        # Left part (Women 48%)
        bar_w_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_x, bar_y, bar_w * 0.48, bar_h)
        bar_w_shape.fill.solid()
        bar_w_shape.fill.fore_color.rgb = COLOR_CYAN
        bar_w_shape.line.fill.background()

        # Right part (Men 52%)
        bar_m_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_x + (bar_w * 0.48), bar_y, bar_w * 0.52, bar_h)
        bar_m_shape.fill.solid()
        bar_m_shape.fill.fore_color.rgb = COLOR_VIOLET
        bar_m_shape.line.fill.background()

        # Center parity marker
        marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_x + (bar_w * 0.5) - Inches(0.02), bar_y - Inches(0.04), Inches(0.04), bar_h + Inches(0.08))
        marker.fill.solid()
        marker.fill.fore_color.rgb = COLOR_WHITE
        marker.line.fill.background()

        # Caption under split bar
        tb_cap = slide.shapes.add_textbox(bar_x, bar_y + Inches(0.38), bar_w, Inches(0.35))
        tf_cap = tb_cap.text_frame
        p_cap = tf_cap.paragraphs[0]
        r_c1 = p_cap.add_run()
        r_c1.text = "Quasi 1 gamer su 2                                "
        r_c1.font.name = FONT_BODY
        r_c1.font.size = Pt(11.5)
        r_c1.font.color.rgb = COLOR_TEXT_MUTED

        r_c2 = p_cap.add_run()
        r_c2.text = "Soglia parità (50%)"
        r_c2.font.name = FONT_BODY
        r_c2.font.size = Pt(11.5)
        r_c2.font.color.rgb = COLOR_TEXT_MUTED

        # Nested Sub-panel: European and Italian stats
        self.add_card(slide, bar_x, y + Inches(2.28), bar_w, Inches(1.22), bg_color=COLOR_SURFACE_SUB, border_color=COLOR_BORDER_MUTED)
        tb_eur = slide.shapes.add_textbox(bar_x + Inches(0.18), y + Inches(2.36), bar_w - Inches(0.36), Inches(1.05))
        tf_eur = tb_eur.text_frame
        tf_eur.word_wrap = True

        p_e1 = tf_eur.paragraphs[0]
        p_e1.text = "• Europa: 47,8% donne · 55 milioni di giocatrici (VGE)"
        p_e1.font.name = FONT_BODY
        p_e1.font.size = Pt(12)
        p_e1.font.bold = True
        p_e1.font.color.rgb = COLOR_TEXT_MAIN

        p_e2 = tf_eur.add_paragraph()
        p_e2.space_before = Pt(6)
        p_e2.text = "• Italia: 41,0% donne · 5,7 milioni su 14M gamer (IIDEA)"
        p_e2.font.name = FONT_BODY
        p_e2.font.size = Pt(12)
        p_e2.font.bold = True
        p_e2.font.color.rgb = COLOR_TEXT_MAIN

        # Right Column: Annual Growth Italy (IIDEA)
        right_x = Inches(6.866)
        self.add_card(slide, right_x, y, col_w, col_h, bg_color=COLOR_SURFACE, border_color=COLOR_BORDER_SUBTLE, left_accent=COLOR_VIOLET)

        tb_r = slide.shapes.add_textbox(right_x + Inches(0.25), y + Inches(0.18), col_w - Inches(0.5), Inches(0.4))
        tf_r = tb_r.text_frame
        p_rk = tf_r.paragraphs[0]
        p_rk.text = "CRESCITA ANNUALE · ITALIA (IIDEA)"
        p_rk.font.name = FONT_MONO
        p_rk.font.size = Pt(11.5)
        p_rk.font.bold = True
        p_rk.font.color.rgb = COLOR_VIOLET

        # Growth Item 1: Women (+14.0%)
        row1_y = y + Inches(0.72)
        tb_grow1 = slide.shapes.add_textbox(right_x + Inches(0.25), row1_y, col_w - Inches(0.5), Inches(0.35))
        tf_g1 = tb_grow1.text_frame
        p_g1 = tf_g1.paragraphs[0]
        r_g1_t = p_g1.add_run()
        r_g1_t.text = "GIOCATRICI DONNE                                       "
        r_g1_t.font.name = FONT_MONO
        r_g1_t.font.bold = True
        r_g1_t.font.size = Pt(12)
        r_g1_t.font.color.rgb = COLOR_TEXT_MAIN

        r_g1_v = p_g1.add_run()
        r_g1_v.text = "+14,0%"
        r_g1_v.font.name = FONT_MONO
        r_g1_v.font.bold = True
        r_g1_v.font.size = Pt(16)
        r_g1_v.font.color.rgb = COLOR_CYAN

        self.add_progress_bar(slide, right_x + Inches(0.25), row1_y + Inches(0.4), col_w - Inches(0.5), Inches(0.22), 1.0, COLOR_CYAN)

        # Growth Item 2: Men (+2.5%)
        row2_y = y + Inches(1.5)
        tb_grow2 = slide.shapes.add_textbox(right_x + Inches(0.25), row2_y, col_w - Inches(0.5), Inches(0.35))
        tf_g2 = tb_grow2.text_frame
        p_g2 = tf_g2.paragraphs[0]
        r_g2_t = p_g2.add_run()
        r_g2_t.text = "GIOCATORI UOMINI                                        "
        r_g2_t.font.name = FONT_MONO
        r_g2_t.font.bold = True
        r_g2_t.font.size = Pt(12)
        r_g2_t.font.color.rgb = COLOR_TEXT_MAIN

        r_g2_v = p_g2.add_run()
        r_g2_v.text = "+2,5%"
        r_g2_v.font.name = FONT_MONO
        r_g2_v.font.bold = True
        r_g2_v.font.size = Pt(16)
        r_g2_v.font.color.rgb = COLOR_VIOLET

        self.add_progress_bar(slide, right_x + Inches(0.25), row2_y + Inches(0.4), col_w - Inches(0.5), Inches(0.22), 0.18, COLOR_VIOLET)

        # Sub-panel description
        self.add_card(slide, right_x + Inches(0.25), y + Inches(2.28), col_w - Inches(0.5), Inches(1.22), bg_color=COLOR_SURFACE_SUB, border_color=COLOR_BORDER_MUTED)
        tb_sub_r = slide.shapes.add_textbox(right_x + Inches(0.43), y + Inches(2.36), col_w - Inches(0.86), Inches(1.05))
        tf_sr = tb_sub_r.text_frame
        tf_sr.word_wrap = True
        p_sr1 = tf_sr.paragraphs[0]
        p_sr1.text = "Dinamica di settore:"
        p_sr1.font.name = FONT_BODY
        p_sr1.font.bold = True
        p_sr1.font.size = Pt(12)
        p_sr1.font.color.rgb = COLOR_TEXT_MAIN

        p_sr2 = tf_sr.add_paragraph()
        p_sr2.space_before = Pt(4)
        p_sr2.text = "L'espansione del pubblico videoludico in Italia è trainata in modo preponderante dall'ingresso attivo della componente femminile."
        p_sr2.font.name = FONT_BODY
        p_sr2.font.size = Pt(11.5)
        p_sr2.font.color.rgb = COLOR_TEXT_MUTED

        # Bottom Quote Bar
        self.add_quote_bar(
            slide,
            left=Inches(0.8), top=Inches(6.16), width=Inches(11.733), height=Inches(0.6),
            regular_text="Il problema non è l'assenza delle donne nel gaming, ma la loro ",
            bold_highlight="invisibilità culturale.",
            accent_color=COLOR_CYAN
        )

    # ----------------------------------------------------------------------------------
    # SLIDE 4: ESPERIENZA ONLINE & TOSSICITÀ (Speaker 2)
    # ----------------------------------------------------------------------------------
    def build_slide_4(self):
        slide = self.create_slide()
        self.add_header(
            slide,
            kicker_text="03 / ESPERIENZA ONLINE",
            kicker_color=COLOR_CORAL,
            title_text="Quando basta una voce per diventare un bersaglio",
            subtitle_text="L'apertura del microfono rompe l'anonimato: l'identità vocale espone a sessismo immediato e ostilità mirata.",
            speaker_text="02 · VALERIO BERNARDI",
            speaker_color=COLOR_VIOLET
        )
        self.add_footer(slide, 4)

        # Top Grid: 2 Large Stat Panels (76% and 48%)
        top_y = Inches(2.22)
        stat_w = Inches(5.666)
        stat_h = Inches(2.18)

        # Panel 1: 76% (ADL Report)
        self.add_card(slide, Inches(0.8), top_y, stat_w, stat_h, bg_color=COLOR_SURFACE, border_color=COLOR_CORAL, left_accent=COLOR_CORAL)
        tb_s1 = slide.shapes.add_textbox(Inches(1.05), top_y + Inches(0.14), stat_w - Inches(0.5), stat_h - Inches(0.25))
        tf_s1 = tb_s1.text_frame
        tf_s1.word_wrap = True

        p_k1 = tf_s1.paragraphs[0]
        p_k1.text = "ADL · HATE IS NO GAME REPORT 2024"
        p_k1.font.name = FONT_MONO
        p_k1.font.size = Pt(11)
        p_k1.font.bold = True
        p_k1.font.color.rgb = COLOR_CORAL

        p_n1 = tf_s1.add_paragraph()
        p_n1.text = "76%"
        p_n1.font.name = FONT_MONO
        p_n1.font.size = Pt(56)
        p_n1.font.bold = True
        p_n1.font.color.rgb = COLOR_CORAL

        p_l1 = tf_s1.add_paragraph()
        r_l1_1 = p_l1.add_run()
        r_l1_1.text = "ha subito "
        r_l1_1.font.name = FONT_BODY
        r_l1_1.font.size = Pt(14)
        r_l1_1.font.color.rgb = COLOR_TEXT_MAIN
        r_l1_2 = p_l1.add_run()
        r_l1_2.text = "molestie nelle partite online competitive"
        r_l1_2.font.name = FONT_BODY
        r_l1_2.font.bold = True
        r_l1_2.font.size = Pt(14)
        r_l1_2.font.color.rgb = COLOR_TEXT_MAIN

        p_sub1 = tf_s1.add_paragraph()
        p_sub1.space_before = Pt(4)
        p_sub1.text = "Insulti verbali, minacce e comportamenti tossici sistematici durante le sessioni."
        p_sub1.font.name = FONT_BODY
        p_sub1.font.size = Pt(11.5)
        p_sub1.font.color.rgb = COLOR_TEXT_MUTED

        # Panel 2: 48% (Target Identity-based)
        self.add_card(slide, Inches(6.866), top_y, stat_w, stat_h, bg_color=COLOR_SURFACE, border_color=COLOR_CORAL, left_accent=COLOR_CORAL)
        tb_s2 = slide.shapes.add_textbox(Inches(7.116), top_y + Inches(0.14), stat_w - Inches(0.5), stat_h - Inches(0.25))
        tf_s2 = tb_s2.text_frame
        tf_s2.word_wrap = True

        p_k2 = tf_s2.paragraphs[0]
        p_k2.text = "PRIMO TARGET IDENTITY-BASED DAL 2019"
        p_k2.font.name = FONT_MONO
        p_k2.font.size = Pt(11)
        p_k2.font.bold = True
        p_k2.font.color.rgb = COLOR_CORAL

        p_n2 = tf_s2.add_paragraph()
        p_n2.text = "48%"
        p_n2.font.name = FONT_MONO
        p_n2.font.size = Pt(56)
        p_n2.font.bold = True
        p_n2.font.color.rgb = COLOR_CORAL

        p_l2 = tf_s2.add_paragraph()
        r_l2_1 = p_l2.add_run()
        r_l2_1.text = "delle giocatrici: "
        r_l2_1.font.name = FONT_BODY
        r_l2_1.font.size = Pt(14)
        r_l2_1.font.color.rgb = COLOR_TEXT_MAIN
        r_l2_2 = p_l2.add_run()
        r_l2_2.text = "molestie mirate specificamente al genere"
        r_l2_2.font.name = FONT_BODY
        r_l2_2.font.bold = True
        r_l2_2.font.size = Pt(14)
        r_l2_2.font.color.rgb = COLOR_TEXT_MAIN

        p_sub2 = tf_s2.add_paragraph()
        p_sub2.space_before = Pt(4)
        p_sub2.text = "Bersaglio primario e più frequente rispetto a qualsiasi altra categoria identitaria."
        p_sub2.font.name = FONT_BODY
        p_sub2.font.size = Pt(11.5)
        p_sub2.font.color.rgb = COLOR_TEXT_MUTED

        # Section Subheader: Come si difendono le giocatrici
        subh_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.55), Inches(8.0), Inches(0.3))
        tf_sh = subh_box.text_frame
        p_sh = tf_sh.paragraphs[0]
        p_sh.text = "STRATEGIE DIFENSIVE ADOTTATE DALLE GIOCATRICI ONLINE"
        p_sh.font.name = FONT_MONO
        p_sh.font.size = Pt(11)
        p_sh.font.bold = True
        p_sh.font.color.rgb = COLOR_TEXT_MUTED

        # Bottom 3 Defense Cards
        defense_items = [
            ("🔇", "Microfono in mute", "Rinuncia forzata alla ", "voice chat", " di squadra."),
            ("🪪", "Nickname neutro", "Mascheramento preventivo dell'", "identità di genere", "."),
            ("🔒", "Solo party chiusi", "Ritiro e auto-esclusione dal ", "matchmaking pubblico", ".")
        ]
        def_w = Inches(3.644)
        def_h = Inches(1.22)
        def_y = Inches(4.9)

        for i, (icon, d_title, pre, hl, post) in enumerate(defense_items):
            dx = Inches(0.8) + i * (def_w + Inches(0.4))
            self.add_card(slide, dx, def_y, def_w, def_h, bg_color=COLOR_SURFACE, border_color=COLOR_BORDER_SUBTLE, left_accent=COLOR_CORAL)

            tb_d = slide.shapes.add_textbox(dx + Inches(0.22), def_y + Inches(0.12), def_w - Inches(0.4), def_h - Inches(0.2))
            tf_d = tb_d.text_frame
            tf_d.word_wrap = True

            p_dt = tf_d.paragraphs[0]
            p_dt.text = f"{icon}  {d_title}"
            p_dt.font.name = FONT_TITLE
            p_dt.font.size = Pt(15)
            p_dt.font.bold = True
            p_dt.font.color.rgb = COLOR_TEXT_MAIN

            p_dd = tf_d.add_paragraph()
            p_dd.space_before = Pt(4)
            r_1 = p_dd.add_run()
            r_1.text = pre
            r_1.font.name = FONT_BODY
            r_1.font.size = Pt(12)
            r_1.font.color.rgb = COLOR_TEXT_MUTED

            r_2 = p_dd.add_run()
            r_2.text = hl
            r_2.font.name = FONT_BODY
            r_2.font.bold = True
            r_2.font.size = Pt(12)
            r_2.font.color.rgb = COLOR_TEXT_MAIN

            r_3 = p_dd.add_run()
            r_3.text = post
            r_3.font.name = FONT_BODY
            r_3.font.size = Pt(12)
            r_3.font.color.rgb = COLOR_TEXT_MUTED

        # Bottom Citation Note
        cite_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.32), Inches(11.733), Inches(0.35))
        tf_c = cite_box.text_frame
        p_c = tf_c.paragraphs[0]
        p_c.text = "FONTI: WELLS ET AL., 2025 (RITIRO DALLA CHAT O ABBANDONO) · GAMERGATE 2014–15 (UN PRECEDENTE ANCORA ATTUALE)"
        p_c.font.name = FONT_MONO
        p_c.font.size = Pt(10.5)
        p_c.font.color.rgb = COLOR_TEXT_MUTED

    # ----------------------------------------------------------------------------------
    # SLIDE 5: CHARACTER DESIGN & RAPPRESENTAZIONE (Speaker 2)
    # ----------------------------------------------------------------------------------
    def build_slide_5(self):
        slide = self.create_slide()
        self.add_header(
            slide,
            kicker_text="04 / CHARACTER DESIGN & RAPPRESENTAZIONE",
            kicker_color=COLOR_VIOLET,
            title_text="Da contorno a Soggetto: L'Evoluzione femminile",
            subtitle_text="Dall'oggettificazione passiva degli esordi alla complessità psicologica e all'identità modulare contemporanea.",
            speaker_text="02 · VALERIO BERNARDI",
            speaker_color=COLOR_VIOLET
        )
        self.add_footer(slide, 5)

        # Timeline Bar Across Top of Nodes
        line_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(2.55), Inches(10.9), Inches(0.04))
        line_bg.fill.solid()
        line_bg.fill.fore_color.rgb = COLOR_BORDER_SUBTLE
        line_bg.line.fill.background()

        # 4 Timeline Era Cards
        eras = [
            (
                "👑", "ARCHETIPO · ANNI '80", COLOR_CORAL,
                "Princess Peach",
                "Donna in pericolo",
                "Personaggio passivo privo di agentività: esiste per essere salvata."
            ),
            (
                "🏹", "1996 · TOMB RAIDER", COLOR_VIOLET,
                "Lara Croft",
                "Male gaze & azione",
                "Costruita per il male gaze visivo, ma autonoma, competente e coraggiosa."
            ),
            (
                "⚔️", "2017–2020 · REALISMO", COLOR_CYAN,
                "Aloy · Senua · Ellie",
                "Complessità & vissuto",
                "Competenza, salute mentale approfondita e corpi realistici non standardizzati."
            ),
            (
                "✨", "OGGI · LIBERTÀ", COLOR_MINT,
                "Identità modulare",
                "Disaccoppiamento",
                "Corpo, voce, estetica e pronomi si separano liberamente (BG3, The Sims 4)."
            )
        ]

        era_w = Inches(2.708)
        era_h = Inches(3.48)
        era_gap = Inches(0.3)
        era_x_start = Inches(0.8)
        era_y = Inches(2.4)

        for i, (icon, tag, color, hero_title, sub_tag, desc) in enumerate(eras):
            x = era_x_start + i * (era_w + era_gap)
            card = self.add_card(slide, x, era_y, era_w, era_h, bg_color=COLOR_SURFACE, border_color=color, left_accent=color)

            tb = slide.shapes.add_textbox(x + Inches(0.18), era_y + Inches(0.18), era_w - Inches(0.36), era_h - Inches(0.3))
            tf = tb.text_frame
            tf.word_wrap = True

            # Era Icon + Tag Pill
            p_tag = tf.paragraphs[0]
            p_tag.text = f"{icon}  {tag}"
            p_tag.font.name = FONT_MONO
            p_tag.font.size = Pt(10)
            p_tag.font.bold = True
            p_tag.font.color.rgb = color

            # Hero Name
            p_name = tf.add_paragraph()
            p_name.space_before = Pt(14)
            p_name.text = hero_title
            p_name.font.name = FONT_TITLE
            p_name.font.size = Pt(18)
            p_name.font.bold = True
            p_name.font.color.rgb = COLOR_TEXT_MAIN

            # Sub-tag (Archetype)
            p_sub = tf.add_paragraph()
            p_sub.space_before = Pt(4)
            p_sub.text = f"[{sub_tag}]"
            p_sub.font.name = FONT_MONO
            p_sub.font.size = Pt(11)
            p_sub.font.bold = True
            p_sub.font.color.rgb = color

            # Description
            p_desc = tf.add_paragraph()
            p_desc.space_before = Pt(12)
            p_desc.text = desc
            p_desc.font.name = FONT_BODY
            p_desc.font.size = Pt(12.5)
            p_desc.font.color.rgb = COLOR_TEXT_MUTED

        # Bottom Quote Bar
        self.add_quote_bar(
            slide,
            left=Inches(0.8), top=Inches(6.16), width=Inches(11.733), height=Inches(0.6),
            regular_text="La figura femminile passa da oggetto della narrazione a ",
            bold_highlight="soggetto della narrazione.",
            accent_color=COLOR_VIOLET
        )

    # ----------------------------------------------------------------------------------
    # SLIDE 6: DIETRO LO SCHERMO: IL LAVORO (Speaker 3)
    # ----------------------------------------------------------------------------------
    def build_slide_6(self):
        slide = self.create_slide()
        self.add_header(
            slide,
            kicker_text="05 / DIETRO LO SCHERMO: IL LAVORO",
            kicker_color=COLOR_MINT,
            title_text="Quasi metà gioca. Solo un quarto sviluppa.",
            subtitle_text="La parità nel consumo non si riflette nella composizione dei team di produzione e ingegneria software.",
            speaker_text="03 · SAMUELE DE SANTIS",
            speaker_color=COLOR_MINT
        )
        self.add_footer(slide, 6)

        # Top Grid: 2 Panels
        top_y = Inches(2.22)
        panel_w = Inches(5.666)
        panel_h = Inches(2.65)

        # Left Panel: GDC State of Industry (23-25%)
        self.add_card(slide, Inches(0.8), top_y, panel_w, panel_h, bg_color=COLOR_SURFACE, border_color=COLOR_CYAN, left_accent=COLOR_CYAN)
        tb_l = slide.shapes.add_textbox(Inches(1.05), top_y + Inches(0.18), panel_w - Inches(0.5), panel_h - Inches(0.3))
        tf_l = tb_l.text_frame
        tf_l.word_wrap = True

        p_lk = tf_l.paragraphs[0]
        p_lk.text = "GDC · STATE OF THE GAME INDUSTRY 2024–2025"
        p_lk.font.name = FONT_MONO
        p_lk.font.size = Pt(11)
        p_lk.font.bold = True
        p_lk.font.color.rgb = COLOR_CYAN

        p_num = tf_l.add_paragraph()
        p_num.space_before = Pt(4)
        p_num.text = "23–25%"
        p_num.font.name = FONT_MONO
        p_num.font.size = Pt(64)
        p_num.font.bold = True
        p_num.font.color.rgb = COLOR_CYAN

        p_lbl = tf_l.add_paragraph()
        p_lbl.text = "donne nella forza lavoro videoludica mondiale"
        p_lbl.font.name = FONT_BODY
        p_lbl.font.size = Pt(15.5)
        p_lbl.font.bold = True
        p_lbl.font.color.rgb = COLOR_TEXT_MAIN

        p_sub = tf_l.add_paragraph()
        p_sub.space_before = Pt(6)
        p_sub.text = "Quota che scende a ≈20% nei ruoli tecnici core e programmazione pura di motori grafici (UE)."
        p_sub.font.name = FONT_BODY
        p_sub.font.size = Pt(12)
        p_sub.font.color.rgb = COLOR_TEXT_MUTED

        # Right Panel: Role & Audience Distribution Bars
        self.add_card(slide, Inches(6.866), top_y, panel_w, panel_h, bg_color=COLOR_SURFACE, border_color=COLOR_BORDER_SUBTLE, left_accent=COLOR_VIOLET)
        tb_r = slide.shapes.add_textbox(Inches(7.116), top_y + Inches(0.18), panel_w - Inches(0.5), Inches(0.4))
        tf_r = tb_r.text_frame
        p_rk = tf_r.paragraphs[0]
        p_rk.text = "DISTRIBUZIONE COMPARATIVA RUOLI E PUBBLICO"
        p_rk.font.name = FONT_MONO
        p_rk.font.size = Pt(11)
        p_rk.font.bold = True
        p_rk.font.color.rgb = COLOR_VIOLET

        # 3 Comparison Progress Bars
        bar_data = [
            ("Pubblico videoludico femminile", "48%", 0.48, COLOR_CYAN, top_y + Inches(0.55)),
            ("Forza lavoro totale nell'industria", "25%", 0.25, COLOR_VIOLET, top_y + Inches(1.25)),
            ("Ruoli tecnici core (Programming)", "≈20%", 0.20, COLOR_CORAL, top_y + Inches(1.95)),
        ]
        for label, val_str, ratio, bar_col, by in bar_data:
            tb_row = slide.shapes.add_textbox(Inches(7.116), by, panel_w - Inches(0.5), Inches(0.3))
            tf_row = tb_row.text_frame
            p_row = tf_row.paragraphs[0]
            r_lbl = p_row.add_run()
            r_lbl.text = f"{label}   "
            r_lbl.font.name = FONT_BODY
            r_lbl.font.size = Pt(11.5)
            r_lbl.font.color.rgb = COLOR_TEXT_MAIN

            r_val = p_row.add_run()
            r_val.text = val_str
            r_val.font.name = FONT_MONO
            r_val.font.bold = True
            r_val.font.size = Pt(12.5)
            r_val.font.color.rgb = bar_col

            self.add_progress_bar(slide, Inches(7.116), by + Inches(0.32), panel_w - Inches(0.5), Inches(0.16), ratio, bar_col)

        # Bottom Grid: Segregazione Orizzontale & Verticale
        bot_y = Inches(5.12)
        bot_h = Inches(1.5)

        # Card 1: Orizzontale
        self.add_card(slide, Inches(0.8), bot_y, panel_w, bot_h, bg_color=COLOR_SURFACE, border_color=COLOR_BORDER_SUBTLE, left_accent=COLOR_VIOLET)
        tb_seg1 = slide.shapes.add_textbox(Inches(1.05), bot_y + Inches(0.15), panel_w - Inches(0.5), bot_h - Inches(0.3))
        tf_s1 = tb_seg1.text_frame
        tf_s1.word_wrap = True
        p_s1_t = tf_s1.paragraphs[0]
        p_s1_t.text = "↔ Segregazione orizzontale"
        p_s1_t.font.name = FONT_TITLE
        p_s1_t.font.size = Pt(17)
        p_s1_t.font.bold = True
        p_s1_t.font.color.rgb = COLOR_TEXT_MAIN

        p_s1_d = tf_s1.add_paragraph()
        p_s1_d.space_before = Pt(6)
        p_s1_d.text = "Uomini e donne concentrati in settori e dipartimenti disuguali: maggiore presenza femminile in QA, narrative, HR e marketing rispetto al software engineering core."
        p_s1_d.font.name = FONT_BODY
        p_s1_d.font.size = Pt(12)
        p_s1_d.font.color.rgb = COLOR_TEXT_MUTED

        # Card 2: Verticale
        self.add_card(slide, Inches(6.866), bot_y, panel_w, bot_h, bg_color=COLOR_SURFACE, border_color=COLOR_BORDER_SUBTLE, left_accent=COLOR_CORAL)
        tb_seg2 = slide.shapes.add_textbox(Inches(7.116), bot_y + Inches(0.15), panel_w - Inches(0.5), bot_h - Inches(0.3))
        tf_s2 = tb_seg2.text_frame
        tf_s2.word_wrap = True
        p_s2_t = tf_s2.paragraphs[0]
        p_s2_t.text = "↥ Segregazione verticale (Soffitto di Cristallo)"
        p_s2_t.font.name = FONT_TITLE
        p_s2_t.font.size = Pt(17)
        p_s2_t.font.bold = True
        p_s2_t.font.color.rgb = COLOR_TEXT_MAIN

        p_s2_d = tf_s2.add_paragraph()
        p_s2_d.space_before = Pt(6)
        p_s2_d.text = "Drammatica rarefazione della presenza femminile salendo nella piramide gerarchica aziendale: barriere strutturali nei ruoli di Direzione Tecnica, Lead e C-Level."
        p_s2_d.font.name = FONT_BODY
        p_s2_d.font.size = Pt(12)
        p_s2_d.font.color.rgb = COLOR_TEXT_MUTED

    # ----------------------------------------------------------------------------------
    # SLIDE 7: POTERE ECONOMICO E COMPETITIVO (Speaker 3)
    # ----------------------------------------------------------------------------------
    def build_slide_7(self):
        slide = self.create_slide()
        self.add_header(
            slide,
            kicker_text="06 / POTERE ECONOMICO E COMPETITIVO",
            kicker_color=COLOR_MINT,
            title_text="Il gap cresce salendo di livello",
            subtitle_text="Dalle asimmetrie retributive alle barriere d'accesso nei circuiti professionistici eSports.",
            speaker_text="03 · SAMUELE DE SANTIS",
            speaker_color=COLOR_MINT
        )
        self.add_footer(slide, 7)

        # Top 3 Stat Panels
        top_y = Inches(2.22)
        col3_w = Inches(3.644)
        col3_h = Inches(2.5)
        gap = Inches(0.4)

        stats = [
            (
                "DIVARIO RETRIBUTIVO", "−24%", COLOR_CORAL,
                "gender pay gap medio negli studi USA",
                "Disparità salariale persistente a parità di livello e ore lavorate nel settore."
            ),
            (
                "SENIORITY & COMPENSI", "68% vs 38%", COLOR_CORAL,
                "oltre 125k$ con 6+ anni di esperienza",
                "Divario netto nell'accesso alle fasce retributive di punta e bonus dirigenziali."
            ),
            (
                "PUBBLICO ESPORTS GLOBALE", "33%", COLOR_CYAN,
                "Spettatrici femminili (Deloitte)",
                "Presenza attiva di giocatrici pro-player fortemente limitata dall'ostilità ambientale."
            ),
        ]

        for i, (kicker, stat_val, col, label, subtext) in enumerate(stats):
            x = Inches(0.8) + i * (col3_w + gap)
            self.add_card(slide, x, top_y, col3_w, col3_h, bg_color=COLOR_SURFACE, border_color=col, left_accent=col)

            tb = slide.shapes.add_textbox(x + Inches(0.2), top_y + Inches(0.15), col3_w - Inches(0.4), col3_h - Inches(0.3))
            tf = tb.text_frame
            tf.word_wrap = True

            p_k = tf.paragraphs[0]
            p_k.text = kicker
            p_k.font.name = FONT_MONO
            p_k.font.size = Pt(10.5)
            p_k.font.bold = True
            p_k.font.color.rgb = col

            p_v = tf.add_paragraph()
            p_v.space_before = Pt(4)
            p_v.text = stat_val
            p_v.font.name = FONT_MONO
            p_v.font.size = Pt(42) if len(stat_val) <= 5 else Pt(34)
            p_v.font.bold = True
            p_v.font.color.rgb = col

            p_l = tf.add_paragraph()
            p_l.space_before = Pt(2)
            p_l.text = label
            p_l.font.name = FONT_BODY
            p_l.font.size = Pt(13)
            p_l.font.bold = True
            p_l.font.color.rgb = COLOR_TEXT_MAIN

            p_s = tf.add_paragraph()
            p_s.space_before = Pt(4)
            p_s.text = subtext
            p_s.font.name = FONT_BODY
            p_s.font.size = Pt(11)
            p_s.font.color.rgb = COLOR_TEXT_MUTED

        # Subheader: Spazi competitivi inclusivi
        sh_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.88), Inches(8.0), Inches(0.3))
        tf_sh = sh_box.text_frame
        p_sh = tf_sh.paragraphs[0]
        p_sh.text = "SPAZI COMPETITIVI INCLUSIVI ED ECOSISTEMI PRO"
        p_sh.font.name = FONT_MONO
        p_sh.font.size = Pt(11)
        p_sh.font.bold = True
        p_sh.font.color.rgb = COLOR_TEXT_MUTED

        # Bottom 2 Esports Cards
        esp_w = Inches(5.666)
        esp_h = Inches(1.4)
        esp_y = Inches(5.22)

        esp_cards = [
            (
                "VCT Game Changers", "Valorant · Riot Games",
                "Circuito professionistico ufficiale e Academy dedicato per creare visibilità, stipendi dignitosi e carriere esports sostenibili per donne e minoranze di genere."
            ),
            (
                "ESL Impact", "Counter-Strike · ESL",
                "Ecosistema competitivo globale dedicato con montepremi dedicati per abbattere le storiche barriere d'ingresso nel panorama competitivo storico di CS2."
            )
        ]

        for i, (title, org, desc) in enumerate(esp_cards):
            ex = Inches(0.8) + i * (esp_w + Inches(0.4))
            self.add_card(slide, ex, esp_y, esp_w, esp_h, bg_color=COLOR_SURFACE, border_color=COLOR_MINT, left_accent=COLOR_MINT)

            tb_e = slide.shapes.add_textbox(ex + Inches(0.22), esp_y + Inches(0.12), esp_w - Inches(0.44), esp_h - Inches(0.24))
            tf_e = tb_e.text_frame
            tf_e.word_wrap = True

            p_et = tf_e.paragraphs[0]
            r_et1 = p_et.add_run()
            r_et1.text = f"◆ {title}  "
            r_et1.font.name = FONT_TITLE
            r_et1.font.size = Pt(15.5)
            r_et1.font.bold = True
            r_et1.font.color.rgb = COLOR_TEXT_MAIN

            r_et2 = p_et.add_run()
            r_et2.text = f"[{org}]"
            r_et2.font.name = FONT_MONO
            r_et2.font.size = Pt(11.5)
            r_et2.font.bold = True
            r_et2.font.color.rgb = COLOR_MINT

            p_ed = tf_e.add_paragraph()
            p_ed.space_before = Pt(6)
            p_ed.text = desc
            p_ed.font.name = FONT_BODY
            p_ed.font.size = Pt(11.5)
            p_ed.font.color.rgb = COLOR_TEXT_MUTED

    # ----------------------------------------------------------------------------------
    # SLIDE 8: DALLA DIVERSITÀ ALL'INCLUSIONE (Speaker 3)
    # ----------------------------------------------------------------------------------
    def build_slide_8(self):
        slide = self.create_slide()
        self.add_header(
            slide,
            kicker_text="07 / DALLA DIVERSITÀ ALL'INCLUSIONE",
            kicker_color=COLOR_MINT,
            title_text="Non basta essere presenti: bisogna poter partecipare",
            subtitle_text="Dalla semplice rappresentazione numerica alla progettazione attiva di spazi e ambienti equi.",
            speaker_text="03 · SAMUELE DE SANTIS",
            speaker_color=COLOR_MINT
        )
        self.add_footer(slide, 8)

        # 3 Pillar Solution Cards
        sol_cols = [
            (
                "01", "DESIGN INCLUSIVO", COLOR_CYAN,
                "Architettura & Personaggi",
                [
                    "Personaggi e trame narrativi liberi da stereotipi riduttivi.",
                    "Sistemi di personalizzazione modulari e slegati da binarismi estetici.",
                    "Meccaniche e stili di gioco che premiano approcci diversificati."
                ]
            ),
            (
                "02", "COMMUNITY SICURE", COLOR_VIOLET,
                "Infrastruttura & Moderazione",
                [
                    "Moderazione proattiva basata su machine learning e reporting chiaro.",
                    "Tolleranza zero per sessismo e molestie verbali in voice chat.",
                    "Algoritmi di matchmaking guidati da indicatori di affidabilità comportamentale."
                ]
            ),
            (
                "03", "OPPORTUNITÀ EQUE", COLOR_MINT,
                "Carriere & Competizione",
                [
                    "Accesso paritario ai ruoli tecnici core e programming (UE).",
                    "Promozione di percorsi di leadership tecnica e trasparenza salariale.",
                    "Supporto ai circuiti professionistici eSports protetti e Academy."
                ]
            )
        ]

        left_start = Inches(0.8)
        col_w = Inches(3.644)
        col_h = Inches(3.72)
        gap = Inches(0.4)
        y = Inches(2.28)

        for i, (num, title, color, sub, bullets) in enumerate(sol_cols):
            x = left_start + i * (col_w + gap)
            self.add_card(slide, x, y, col_w, col_h, bg_color=COLOR_SURFACE, border_color=color, left_accent=color)

            tb = slide.shapes.add_textbox(x + Inches(0.22), y + Inches(0.18), col_w - Inches(0.44), col_h - Inches(0.36))
            tf = tb.text_frame
            tf.word_wrap = True

            p_n = tf.paragraphs[0]
            p_n.text = num
            p_n.font.name = FONT_MONO
            p_n.font.size = Pt(46)
            p_n.font.bold = True
            p_n.font.color.rgb = color

            p_t = tf.add_paragraph()
            p_t.text = title
            p_t.font.name = FONT_TITLE
            p_t.font.size = Pt(19)
            p_t.font.bold = True
            p_t.font.color.rgb = COLOR_TEXT_MAIN

            p_s = tf.add_paragraph()
            p_s.text = f"[{sub}]"
            p_s.font.name = FONT_MONO
            p_s.font.size = Pt(11)
            p_s.font.bold = True
            p_s.font.color.rgb = color
            p_s.space_after = Pt(12)

            for b in bullets:
                p_b = tf.add_paragraph()
                p_b.space_before = Pt(6)
                p_b.text = f"• {b}"
                p_b.font.name = FONT_BODY
                p_b.font.size = Pt(12)
                p_b.font.color.rgb = COLOR_TEXT_MUTED

        # Bottom Quote Bar
        self.add_quote_bar(
            slide,
            left=Inches(0.8), top=Inches(6.16), width=Inches(11.733), height=Inches(0.6),
            regular_text="L'inclusione non è un'aggiunta al prodotto: è un ",
            bold_highlight="requisito fondamentale di progettazione.",
            accent_color=COLOR_MINT
        )

    # ----------------------------------------------------------------------------------
    # SLIDE 9: SINTESI E CONCLUSIONI (Speaker 3)
    # ----------------------------------------------------------------------------------
    def build_slide_9(self):
        slide = self.create_slide()
        self.add_header(
            slide,
            kicker_text="08 / SINTESI E CONCLUSIONI",
            kicker_color=COLOR_MINT,
            title_text="La presenza non è ancora piena inclusione",
            subtitle_text="Sintesi del percorso: dai dati di mercato alla sfida culturale e ingegneristica per il futuro.",
            speaker_text="03 · SAMUELE DE SANTIS",
            speaker_color=COLOR_MINT
        )
        self.add_footer(slide, 9)

        # 2 Main Columns
        y = Inches(2.28)
        left_w = Inches(6.8)
        right_w = Inches(4.533)
        h = Inches(3.9)

        # Left Column: Big Claim, Vision & Quote
        tb_l = slide.shapes.add_textbox(Inches(0.8), y, left_w, h)
        tf_l = tb_l.text_frame
        tf_l.word_wrap = True

        p_claim = tf_l.paragraphs[0]
        r_c1 = p_claim.add_run()
        r_c1.text = "Le donne nel gaming "
        r_c1.font.name = FONT_TITLE
        r_c1.font.size = Pt(36)
        r_c1.font.bold = True
        r_c1.font.color.rgb = COLOR_TEXT_MAIN

        r_c2 = p_claim.add_run()
        r_c2.text = "ci sono già."
        r_c2.font.name = FONT_TITLE
        r_c2.font.size = Pt(36)
        r_c2.font.bold = True
        r_c2.font.color.rgb = COLOR_CYAN

        p_desc = tf_l.add_paragraph()
        p_desc.space_before = Pt(14)
        p_desc.text = "La sfida non consiste nell'attrarre un pubblico che già partecipa quotidianamente con passione e competenza, ma nel trasformare l'intero ecosistema tecnologico affinché tutte e tutti possano viverlo, lavorarci e competere alle medesime condizioni di sicurezza, rispetto e pari opportunità."
        p_desc.font.name = FONT_BODY
        p_desc.font.size = Pt(15)
        p_desc.font.color.rgb = COLOR_TEXT_MUTED

        # Quote box embedded
        quote_y = y + Inches(2.15)
        self.add_quote_bar(
            slide,
            left=Inches(0.8), top=quote_y, width=left_w, height=Inches(0.7),
            regular_text="L’inclusione parte dalla ",
            bold_highlight="valorizzazione dell’unicità della persona.",
            accent_color=COLOR_VIOLET
        )

        # Bottom Tor Vergata Badge
        tag_badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y + Inches(3.1), Inches(4.8), Inches(0.36))
        tag_badge.fill.solid()
        tag_badge.fill.fore_color.rgb = COLOR_SURFACE
        tag_badge.line.color.rgb = COLOR_VIOLET
        tag_badge.line.width = Pt(1)
        tf_tb = tag_badge.text_frame
        tf_tb.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf_tb.margin_left = Inches(0.12)
        p_tb = tf_tb.paragraphs[0]
        p_tb.text = "● GENDER & INCLUSION · TOR VERGATA · 2025–2026"
        p_tb.font.name = FONT_MONO
        p_tb.font.size = Pt(10.5)
        p_tb.font.bold = True
        p_tb.font.color.rgb = COLOR_VIOLET

        # Right Column: Academic Sources Panel
        rx = Inches(8.0)
        self.add_card(slide, rx, y, right_w, h, bg_color=COLOR_SURFACE, border_color=COLOR_CYAN, left_accent=COLOR_CYAN)

        tb_r = slide.shapes.add_textbox(rx + Inches(0.25), y + Inches(0.16), right_w - Inches(0.5), h - Inches(0.32))
        tf_r = tb_r.text_frame
        tf_r.word_wrap = True

        p_src_h = tf_r.paragraphs[0]
        p_src_h.text = "⚡ FONTI PRINCIPALI & ACCADEMICHE"
        p_src_h.font.name = FONT_MONO
        p_src_h.font.size = Pt(12)
        p_src_h.font.bold = True
        p_src_h.font.color.rgb = COLOR_CYAN

        sources = [
            ("ESA (2025)", "Global Power of Play Report"),
            ("ADL (2024)", "Hate is No Game Report"),
            ("IIDEA (2025)", "I videogiochi in Italia nel 2024"),
            ("GDC (2025)", "State of the Game Industry Report"),
            ("VGE (2024)", "Key Facts Report (GameTrack)"),
            ("Wells et al. (2025)", "Frontiers in Psychology"),
            ("Botto (2025)", "AboutGender Journal"),
        ]

        for author, work in sources:
            p_s = tf_r.add_paragraph()
            p_s.space_before = Pt(7)
            r_auth = p_s.add_run()
            r_auth.text = f"• {author} — "
            r_auth.font.name = FONT_MONO
            r_auth.font.bold = True
            r_auth.font.size = Pt(11)
            r_auth.font.color.rgb = COLOR_CYAN

            r_work = p_s.add_run()
            r_work.text = work
            r_work.font.name = FONT_BODY
            r_work.font.size = Pt(11)
            r_work.font.color.rgb = COLOR_TEXT_MUTED

        # Bottom Dossier Academic Note
        dossier_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.35), Inches(11.733), Inches(0.35))
        tf_doss = dossier_box.text_frame
        p_doss = tf_doss.paragraphs[0]
        p_doss.text = "DOSSIER ACCADEMICO COMPLETO CON BIBLIOGRAFIA IN FORMATO APA DISPONIBILE NEGLI ATTI DEL CORSO"
        p_doss.font.name = FONT_MONO
        p_doss.font.size = Pt(10.5)
        p_doss.font.color.rgb = COLOR_TEXT_MUTED

    # ----------------------------------------------------------------------------------
    # BUILD ALL 9 SLIDES & SAVE
    # ----------------------------------------------------------------------------------
    def build_all(self, output_paths):
        print("[*] Building Slide 1: Cover & Team...")
        self.build_slide_1()

        print("[*] Building Slide 2: Prospettiva Informatica (Luca Gugliotta)...")
        self.build_slide_2()

        print("[*] Building Slide 3: Demografia e Mercato (Luca Gugliotta)...")
        self.build_slide_3()

        print("[*] Building Slide 4: Esperienza Online & Molestie (Valerio Bernardi)...")
        self.build_slide_4()

        print("[*] Building Slide 5: Character Design & Rappresentazione (Valerio Bernardi)...")
        self.build_slide_5()

        print("[*] Building Slide 6: Dietro lo Schermo: Il Lavoro (Samuele De Santis)...")
        self.build_slide_6()

        print("[*] Building Slide 7: Potere Economico ed Esports (Samuele De Santis)...")
        self.build_slide_7()

        print("[*] Building Slide 8: Dalla Diversità all'Inclusione (Samuele De Santis)...")
        self.build_slide_8()

        print("[*] Building Slide 9: Sintesi e Conclusioni (Samuele De Santis)...")
        self.build_slide_9()

        for path in output_paths:
            os.makedirs(os.path.dirname(os.path.abspath(path)), exist_ok=True)
            self.prs.save(path)
            size_kb = os.path.getsize(path) / 1024
            print(f"[+] Successfully saved presentation: {path} ({size_kb:.1f} KB)")


def main():
    print("=" * 70)
    print("GENERAZIONE POWERPOINT 100% EDITABILE · DONNE, GAMING E INCLUSIONE")
    print("=" * 70)

    builder = EditableDeckBuilder()

    # Paths to save
    script_dir = os.path.dirname(os.path.abspath(__file__))
    project_folder_path = os.path.join(
        script_dir,
        "content", "UNI", "ANNO 3", "GENDER INCLUSION", "PROGETTO_GAMING",
        "donne-gaming-inclusione-modificabile.pptx"
    )
    root_path = os.path.join(script_dir, "donne-gaming-inclusione-modificabile.pptx")

    # If already running inside PROGETTO_GAMING
    if os.path.basename(script_dir) == "PROGETTO_GAMING":
        project_folder_path = os.path.join(script_dir, "donne-gaming-inclusione-modificabile.pptx")
        root_path = os.path.abspath(os.path.join(script_dir, "..", "..", "..", "..", "donne-gaming-inclusione-modificabile.pptx"))

    output_paths = [project_folder_path, root_path]
    # Remove duplicates
    output_paths = list(dict.fromkeys(output_paths))

    builder.build_all(output_paths)
    print("=" * 70)
    print("COMPLETATO CON SUCCESSO! Deck 100% vettoriale ed editabile pronto.")
    print("=" * 70)


if __name__ == "__main__":
    main()
